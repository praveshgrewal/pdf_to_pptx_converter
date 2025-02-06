import os
import streamlit as st
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches

def pdf_to_pptx(pdf_path, pptx_path):
    try:
        images = convert_from_path(pdf_path)
        prs = Presentation()
        
        for img in images:
            slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide
            img_path = "temp_img.png"
            img.save(img_path, "PNG")
            left = top = Inches(0.5)
            slide.shapes.add_picture(img_path, left, top, width=Inches(9))
            os.remove(img_path)
        
        prs.save(pptx_path)
        return pptx_path
    except Exception as e:
        return str(e)

st.title("PDF to PPTX Converter")

uploaded_file = st.file_uploader("Upload a PDF file", type=["pdf"])

if uploaded_file is not None:
    pdf_path = "uploaded.pdf"
    with open(pdf_path, "wb") as f:
        f.write(uploaded_file.getbuffer())
    
    pptx_path = "converted.pptx"
    result = pdf_to_pptx(pdf_path, pptx_path)
    
    if os.path.exists(pptx_path):
        with open(pptx_path, "rb") as f:
            st.download_button("Download PPTX", f, file_name="converted.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
    else:
        st.error(f"Error: {result}")

